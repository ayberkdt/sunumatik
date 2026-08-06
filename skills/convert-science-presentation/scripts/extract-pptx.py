#!/usr/bin/env python3
"""Extract a conservative PPTX inventory without modifying the source deck."""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path


def stable_id(slide_number: int, shape_number: int, kind: str) -> str:
    raw = f"{slide_number}:{shape_number}:{kind}".encode("utf-8")
    return "pptx-" + hashlib.sha1(raw).hexdigest()[:12]


def shape_kind(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "shape_type", None) == 13:
        return "picture"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return str(getattr(getattr(shape, "shape_type", None), "name", "unknown")).lower()


def extract(source: Path, assets_dir: Path | None = None) -> dict:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required. Install it in the project environment, then rerun."
        ) from exc

    deck = Presentation(str(source))
    if assets_dir:
        assets_dir.mkdir(parents=True, exist_ok=True)
    result = {
        "schemaVersion": 1,
        "source": str(source.resolve()),
        "slideSizeEmu": {"width": deck.slide_width, "height": deck.slide_height},
        "slides": [],
        "unsupported": [],
    }
    if assets_dir:
        result["assetsDir"] = str(assets_dir.resolve())

    for slide_number, slide in enumerate(deck.slides, start=1):
        slide_item = {"id": f"slide-{slide_number:03d}", "number": slide_number, "shapes": []}
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_item["notes"] = notes
        except Exception:
            pass

        for shape_number, shape in enumerate(slide.shapes, start=1):
            kind = shape_kind(shape)
            item = {
                "id": stable_id(slide_number, shape_number, kind),
                "name": shape.name,
                "kind": kind,
                "boxEmu": {
                    "x": shape.left,
                    "y": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                },
            }
            if getattr(shape, "has_text_frame", False):
                item["text"] = shape.text
                item["paragraphs"] = [
                    {
                        "level": paragraph.level,
                        "text": paragraph.text,
                        "runs": [
                            {
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                                "font": run.font.name,
                                "sizePt": run.font.size.pt if run.font.size else None,
                            }
                            for run in paragraph.runs
                        ],
                    }
                    for paragraph in shape.text_frame.paragraphs
                ]
            if getattr(shape, "has_table", False):
                item["cells"] = [[cell.text for cell in row.cells] for row in shape.table.rows]
            if getattr(shape, "has_chart", False):
                item["chartType"] = str(shape.chart.chart_type)
                result["unsupported"].append(
                    {"slideId": slide_item["id"], "shapeId": item["id"], "reason": "chart-data-review"}
                )
            if kind == "picture":
                blob = shape.image.blob
                digest = hashlib.sha256(blob).hexdigest()
                extension = shape.image.ext or "bin"
                asset = {
                    "sha256": digest,
                    "bytes": len(blob),
                    "contentType": shape.image.content_type,
                }
                if assets_dir:
                    filename = f"slide-{slide_number:03d}-shape-{shape_number:03d}-{digest[:12]}.{extension}"
                    destination = assets_dir / filename
                    if not destination.exists():
                        destination.write_bytes(blob)
                    asset["path"] = str(destination.resolve())
                else:
                    asset["embedded"] = True
                item["asset"] = asset
            if kind not in {"text", "table", "picture", "chart", "placeholder"}:
                result["unsupported"].append(
                    {"slideId": slide_item["id"], "shapeId": item["id"], "reason": kind}
                )
            slide_item["shapes"].append(item)
        result["slides"].append(slide_item)
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source", type=Path, help="Source .pptx file")
    parser.add_argument("--output", "-o", type=Path, help="Write JSON here; defaults to stdout")
    parser.add_argument("--assets-dir", type=Path, help="Extract embedded pictures into this folder")
    args = parser.parse_args()
    if not args.source.is_file() or args.source.suffix.lower() != ".pptx":
        parser.error("source must be an existing .pptx file")
    payload = json.dumps(extract(args.source, args.assets_dir), ensure_ascii=False, indent=2)
    if args.output:
        args.output.write_text(payload + "\n", encoding="utf-8")
    else:
        sys.stdout.write(payload + "\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
